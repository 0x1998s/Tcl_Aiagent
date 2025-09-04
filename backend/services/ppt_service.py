"""
PPT自动生成服务
支持数据报告的PPT自动生成和模板管理
"""

import os
import json
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import base64
from io import BytesIO

from utils.logger import get_logger
from core.config import Settings

logger = get_logger(__name__)


class PPTService:
    """PPT自动生成服务"""
    
    def __init__(self, settings: Settings):
        self.settings = settings
        self.template_dir = Path("./templates/ppt")
        self.output_dir = Path(settings.REPORT_OUTPUT_DIR)
        self.charts_dir = Path("./temp/charts")
        
        # 确保目录存在
        self.template_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.charts_dir.mkdir(parents=True, exist_ok=True)
        
    async def initialize(self):
        """初始化PPT服务"""
        logger.info("初始化PPT生成服务...")
        
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx未安装，PPT功能将受限")
            return
        
        # 创建默认模板
        await self._create_default_templates()
        
        logger.info("PPT生成服务初始化完成")
    
    async def _create_default_templates(self):
        """创建默认PPT模板"""
        
        templates = {
            "business_report": {
                "name": "商业报告模板",
                "slides": [
                    {
                        "type": "title",
                        "title": "{{report_title}}",
                        "subtitle": "{{report_period}}"
                    },
                    {
                        "type": "overview",
                        "title": "核心指标概览",
                        "content": "{{key_metrics}}"
                    },
                    {
                        "type": "chart",
                        "title": "趋势分析",
                        "chart_type": "line",
                        "data_source": "{{trend_data}}"
                    },
                    {
                        "type": "insights",
                        "title": "关键洞察",
                        "content": "{{insights}}"
                    },
                    {
                        "type": "recommendations",
                        "title": "行动建议",
                        "content": "{{recommendations}}"
                    }
                ]
            },
            "data_analysis": {
                "name": "数据分析模板",
                "slides": [
                    {
                        "type": "title",
                        "title": "数据分析报告",
                        "subtitle": "{{analysis_date}}"
                    },
                    {
                        "type": "methodology",
                        "title": "分析方法",
                        "content": "{{methodology}}"
                    },
                    {
                        "type": "findings",
                        "title": "主要发现",
                        "content": "{{findings}}"
                    },
                    {
                        "type": "chart",
                        "title": "数据可视化",
                        "chart_type": "multiple",
                        "data_source": "{{charts_data}}"
                    }
                ]
            }
        }
        
        # 保存模板配置
        for template_name, template_config in templates.items():
            template_path = self.template_dir / f"{template_name}.json"
            with open(template_path, "w", encoding="utf-8") as f:
                json.dump(template_config, f, ensure_ascii=False, indent=2)
    
    async def generate_ppt(
        self,
        data: Dict[str, Any],
        template_name: str = "business_report",
        output_filename: Optional[str] = None
    ) -> Dict[str, Any]:
        """生成PPT报告"""
        
        if not PPTX_AVAILABLE:
            return await self._generate_html_report(data, template_name, output_filename)
        
        try:
            logger.info(f"开始生成PPT: {template_name}")
            
            # 加载模板
            template = await self._load_template(template_name)
            if not template:
                raise ValueError(f"模板不存在: {template_name}")
            
            # 创建PPT
            prs = Presentation()
            
            # 生成幻灯片
            for slide_config in template["slides"]:
                await self._create_slide(prs, slide_config, data)
            
            # 保存PPT
            if not output_filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_filename = f"report_{timestamp}.pptx"
            
            output_path = self.output_dir / output_filename
            prs.save(str(output_path))
            
            logger.info(f"PPT生成完成: {output_path}")
            
            return {
                "status": "success",
                "file_path": str(output_path),
                "file_size": output_path.stat().st_size,
                "slide_count": len(prs.slides),
                "template_used": template_name
            }
            
        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}")
            return {
                "status": "error",
                "error": str(e)
            }
    
    async def _load_template(self, template_name: str) -> Optional[Dict[str, Any]]:
        """加载PPT模板"""
        template_path = self.template_dir / f"{template_name}.json"
        
        if not template_path.exists():
            return None
        
        try:
            with open(template_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"加载模板失败: {str(e)}")
            return None
    
    async def _create_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建单个幻灯片"""
        
        slide_type = slide_config.get("type", "content")
        
        if slide_type == "title":
            await self._create_title_slide(prs, slide_config, data)
        elif slide_type == "overview":
            await self._create_overview_slide(prs, slide_config, data)
        elif slide_type == "chart":
            await self._create_chart_slide(prs, slide_config, data)
        elif slide_type == "insights":
            await self._create_insights_slide(prs, slide_config, data)
        elif slide_type == "recommendations":
            await self._create_recommendations_slide(prs, slide_config, data)
        else:
            await self._create_content_slide(prs, slide_config, data)
    
    async def _create_title_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建标题幻灯片"""
        
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self._replace_placeholders(slide_config.get("title", ""), data)
        subtitle.text = self._replace_placeholders(slide_config.get("subtitle", ""), data)
    
    async def _create_overview_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建概览幻灯片"""
        
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = self._replace_placeholders(slide_config.get("title", ""), data)
        
        # 添加核心指标
        if "key_metrics" in data:
            metrics_text = self._format_key_metrics(data["key_metrics"])
            content.text = metrics_text
    
    async def _create_chart_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建图表幻灯片"""
        
        slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_shape.text_frame
        title_frame.text = self._replace_placeholders(slide_config.get("title", ""), data)
        
        # 生成图表
        chart_type = slide_config.get("chart_type", "line")
        chart_data = data.get("charts_data", {})
        
        if chart_data:
            chart_path = await self._generate_chart(chart_data, chart_type)
            if chart_path:
                # 添加图表图片
                slide.shapes.add_picture(
                    str(chart_path),
                    Inches(1), Inches(2), Inches(8), Inches(5)
                )
    
    async def _create_insights_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建洞察幻灯片"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = self._replace_placeholders(slide_config.get("title", ""), data)
        
        # 添加洞察内容
        if "insights" in data:
            insights_text = self._format_insights(data["insights"])
            content.text = insights_text
    
    async def _create_recommendations_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建建议幻灯片"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = self._replace_placeholders(slide_config.get("title", ""), data)
        
        # 添加建议内容
        if "recommendations" in data:
            recommendations_text = self._format_recommendations(data["recommendations"])
            content.text = recommendations_text
    
    async def _create_content_slide(
        self, 
        prs: Presentation, 
        slide_config: Dict[str, Any], 
        data: Dict[str, Any]
    ):
        """创建通用内容幻灯片"""
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = self._replace_placeholders(slide_config.get("title", ""), data)
        content.text = self._replace_placeholders(slide_config.get("content", ""), data)
    
    def _replace_placeholders(self, text: str, data: Dict[str, Any]) -> str:
        """替换占位符"""
        
        for key, value in data.items():
            placeholder = f"{{{{{key}}}}}"
            if placeholder in text:
                text = text.replace(placeholder, str(value))
        
        return text
    
    def _format_key_metrics(self, metrics: Dict[str, Any]) -> str:
        """格式化关键指标"""
        
        formatted_text = ""
        for metric_name, metric_value in metrics.items():
            if isinstance(metric_value, dict):
                value = metric_value.get("value", "N/A")
                change = metric_value.get("change", "")
                formatted_text += f"• {metric_name}: {value} {change}\n"
            else:
                formatted_text += f"• {metric_name}: {metric_value}\n"
        
        return formatted_text
    
    def _format_insights(self, insights: List[str]) -> str:
        """格式化洞察内容"""
        
        formatted_text = ""
        for i, insight in enumerate(insights, 1):
            formatted_text += f"{i}. {insight}\n\n"
        
        return formatted_text
    
    def _format_recommendations(self, recommendations: List[str]) -> str:
        """格式化建议内容"""
        
        formatted_text = ""
        for i, recommendation in enumerate(recommendations, 1):
            formatted_text += f"建议{i}: {recommendation}\n\n"
        
        return formatted_text
    
    async def _generate_chart(
        self, 
        chart_data: Dict[str, Any], 
        chart_type: str
    ) -> Optional[Path]:
        """生成图表图片"""
        
        try:
            plt.style.use('default')
            fig, ax = plt.subplots(figsize=(10, 6))
            
            if chart_type == "line":
                await self._create_line_chart(ax, chart_data)
            elif chart_type == "bar":
                await self._create_bar_chart(ax, chart_data)
            elif chart_type == "pie":
                await self._create_pie_chart(ax, chart_data)
            else:
                await self._create_line_chart(ax, chart_data)  # 默认线图
            
            # 保存图表
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            chart_path = self.charts_dir / f"chart_{timestamp}.png"
            
            plt.tight_layout()
            plt.savefig(str(chart_path), dpi=300, bbox_inches='tight')
            plt.close()
            
            return chart_path
            
        except Exception as e:
            logger.error(f"生成图表失败: {str(e)}")
            return None
    
    async def _create_line_chart(self, ax, chart_data: Dict[str, Any]):
        """创建线图"""
        
        x_data = chart_data.get("x", [])
        y_data = chart_data.get("y", [])
        title = chart_data.get("title", "趋势图")
        
        ax.plot(x_data, y_data, marker='o', linewidth=2, markersize=6)
        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.set_xlabel(chart_data.get("x_label", ""), fontsize=12)
        ax.set_ylabel(chart_data.get("y_label", ""), fontsize=12)
        ax.grid(True, alpha=0.3)
    
    async def _create_bar_chart(self, ax, chart_data: Dict[str, Any]):
        """创建柱状图"""
        
        x_data = chart_data.get("x", [])
        y_data = chart_data.get("y", [])
        title = chart_data.get("title", "柱状图")
        
        ax.bar(x_data, y_data, alpha=0.8)
        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.set_xlabel(chart_data.get("x_label", ""), fontsize=12)
        ax.set_ylabel(chart_data.get("y_label", ""), fontsize=12)
    
    async def _create_pie_chart(self, ax, chart_data: Dict[str, Any]):
        """创建饼图"""
        
        labels = chart_data.get("labels", [])
        sizes = chart_data.get("sizes", [])
        title = chart_data.get("title", "饼图")
        
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90)
        ax.set_title(title, fontsize=16, fontweight='bold')
    
    async def _generate_html_report(
        self,
        data: Dict[str, Any],
        template_name: str,
        output_filename: Optional[str] = None
    ) -> Dict[str, Any]:
        """生成HTML报告（当PPT库不可用时的备选方案）"""
        
        try:
            if not output_filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_filename = f"report_{timestamp}.html"
            
            html_content = self._generate_html_content(data, template_name)
            
            output_path = self.output_dir / output_filename
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(html_content)
            
            return {
                "status": "success",
                "file_path": str(output_path),
                "file_size": output_path.stat().st_size,
                "format": "html",
                "template_used": template_name
            }
            
        except Exception as e:
            logger.error(f"HTML报告生成失败: {str(e)}")
            return {
                "status": "error",
                "error": str(e)
            }
    
    def _generate_html_content(self, data: Dict[str, Any], template_name: str) -> str:
        """生成HTML内容"""
        
        html_template = """
        <!DOCTYPE html>
        <html lang="zh-CN">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{title}</title>
            <style>
                body {{ font-family: 'Microsoft YaHei', Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
                h2 {{ color: #34495e; margin-top: 30px; }}
                .metric {{ background: #f8f9fa; padding: 15px; margin: 10px 0; border-left: 4px solid #3498db; }}
                .insight {{ background: #e8f5e8; padding: 15px; margin: 10px 0; border-radius: 5px; }}
                .recommendation {{ background: #fff3cd; padding: 15px; margin: 10px 0; border-radius: 5px; }}
            </style>
        </head>
        <body>
            <h1>{title}</h1>
            <p><strong>生成时间:</strong> {timestamp}</p>
            
            <h2>核心指标</h2>
            {metrics_html}
            
            <h2>关键洞察</h2>
            {insights_html}
            
            <h2>行动建议</h2>
            {recommendations_html}
        </body>
        </html>
        """
        
        # 格式化数据
        title = data.get("report_title", "数据分析报告")
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        metrics_html = ""
        if "key_metrics" in data:
            for metric_name, metric_value in data["key_metrics"].items():
                metrics_html += f'<div class="metric"><strong>{metric_name}:</strong> {metric_value}</div>'
        
        insights_html = ""
        if "insights" in data:
            for insight in data["insights"]:
                insights_html += f'<div class="insight">{insight}</div>'
        
        recommendations_html = ""
        if "recommendations" in data:
            for recommendation in data["recommendations"]:
                recommendations_html += f'<div class="recommendation">{recommendation}</div>'
        
        return html_template.format(
            title=title,
            timestamp=timestamp,
            metrics_html=metrics_html,
            insights_html=insights_html,
            recommendations_html=recommendations_html
        )
    
    async def get_available_templates(self) -> List[Dict[str, Any]]:
        """获取可用的PPT模板"""
        
        templates = []
        
        for template_file in self.template_dir.glob("*.json"):
            try:
                with open(template_file, "r", encoding="utf-8") as f:
                    template_config = json.load(f)
                    
                templates.append({
                    "name": template_file.stem,
                    "display_name": template_config.get("name", template_file.stem),
                    "slide_count": len(template_config.get("slides", [])),
                    "description": template_config.get("description", "")
                })
                
            except Exception as e:
                logger.error(f"加载模板失败 {template_file}: {str(e)}")
        
        return templates
